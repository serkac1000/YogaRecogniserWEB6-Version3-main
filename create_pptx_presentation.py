
#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_pose_recognition_presentation():
    # Create a presentation object
    prs = Presentation()
    
    # Set slide layout
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    bullet_slide_layout = prs.slide_layouts[1]  # Title and content
    
    # Define colors
    primary_color = RGBColor(0, 102, 204)  # Blue
    secondary_color = RGBColor(255, 153, 51)  # Orange
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Pose Recognition Web App"
    subtitle.text = "AI-Powered Real-Time Pose Detection with Local Model Support\n\nReplit Development Team\nJuly 09, 2025"
    
    # Slide 2: What is Pose Recognition App?
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What is Pose Recognition App?"
    tf = content.text_frame
    tf.text = "Real-time pose detection using webcam input"
    
    p = tf.add_paragraph()
    p.text = "Offline functionality - no internet required after setup"
    p = tf.add_paragraph()
    p.text = "Customizable settings for 1-7 poses"
    p = tf.add_paragraph()
    p.text = "Local model support with Teachable Machine integration"
    p = tf.add_paragraph()
    p.text = "Distance calibration for optimal positioning"
    p = tf.add_paragraph()
    p.text = "Audio feedback and visual guidance"
    
    # Slide 3: Core Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Core Features"
    tf = content.text_frame
    tf.text = "Pose Recognition System:"
    
    p = tf.add_paragraph()
    p.text = "• Support for 1-7 customizable poses"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Real-time confidence scoring (30-90% threshold)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Visual pose comparison with reference images"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Skeletal tracking with 17 keypoints"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Customization Options:"
    p = tf.add_paragraph()
    p.text = "• Adjustable recognition delay (1-10 seconds)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Custom pose names (editable labels)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Audio beep toggle for success feedback"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Distance calibration with real-time guidance"
    p.level = 1
    
    # Slide 4: GUI Settings Interface
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "GUI: Settings Interface"
    tf = content.text_frame
    tf.text = "Key Features:"
    
    p = tf.add_paragraph()
    p.text = "• Local model file uploads (model.json, metadata.json, weights.bin)"
    p = tf.add_paragraph()
    p.text = "• Pose selection checkboxes for 1-7 poses"
    p = tf.add_paragraph()
    p.text = "• Reference image uploads for each pose"
    p = tf.add_paragraph()
    p.text = "• Audio, delay, and accuracy threshold controls"
    p = tf.add_paragraph()
    p.text = "• Distance calibration guide"
    
    # Try to add image if it exists
    if os.path.exists("attached_assets/GUI1_1752049448296.png"):
        slide.shapes.add_picture("attached_assets/GUI1_1752049448296.png", 
                               Inches(1), Inches(3), Inches(8), Inches(4))
    
    # Slide 5: GUI Recognition Interface
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "GUI: Recognition Interface"
    tf = content.text_frame
    tf.text = "Real-time Features:"
    
    p = tf.add_paragraph()
    p.text = "• Live webcam feed with pose detection"
    p = tf.add_paragraph()
    p.text = "• Current vs Expected pose display"
    p = tf.add_paragraph()
    p.text = "• Confidence bar with color coding (green = correct, red = incorrect)"
    p = tf.add_paragraph()
    p.text = "• Reference image comparison"
    p = tf.add_paragraph()
    p.text = "• Refresh and back to settings controls"
    
    # Try to add image if it exists
    if os.path.exists("attached_assets/GUI2_1752049449853.png"):
        slide.shapes.add_picture("attached_assets/GUI2_1752049449853.png", 
                               Inches(1), Inches(3), Inches(8), Inches(4))
    
    # Slide 6: Technical Specifications
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Specifications"
    tf = content.text_frame
    tf.text = "AI Framework:"
    
    p = tf.add_paragraph()
    p.text = "• TensorFlow.js with Teachable Machine models"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• PoseNet for 17-keypoint skeletal tracking"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Local storage with IndexedDB for offline use"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Compatibility:"
    p = tf.add_paragraph()
    p.text = "• Modern browsers with WebRTC support"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Optimized for Windows compatibility"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Automatic image compression for storage"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Multiple camera resolution fallback"
    p.level = 1
    
    # Slide 7: Easy Setup Process
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Easy Setup Process"
    tf = content.text_frame
    tf.text = "1. Train Your Model: Use Teachable Machine to create pose model"
    
    p = tf.add_paragraph()
    p.text = "2. Upload Files: Import model.json, metadata.json, weights.bin"
    p = tf.add_paragraph()
    p.text = "3. Configure Poses: Select 1-7 poses and upload reference images"
    p = tf.add_paragraph()
    p.text = "4. Adjust Settings: Set audio, delay, and accuracy preferences"
    p = tf.add_paragraph()
    p.text = "5. Start Recognition: Begin real-time pose detection"
    
    # Add note box
    p = tf.add_paragraph()
    p.text = "\nFully Offline: Works completely without internet after initial setup"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)
    
    # Slide 8: Advanced Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Advanced Features"
    tf = content.text_frame
    tf.text = "Distance Calibration:"
    
    p = tf.add_paragraph()
    p.text = "• Auto-detection of user distance from camera"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Real-time feedback for optimal positioning (3-4 feet)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Visual cues: Green = perfect, Red = adjust distance"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Management:"
    p = tf.add_paragraph()
    p.text = "• Save all settings and model files locally"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Clear memory function with confirmation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Auto-save for pose selections and custom names"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Persistent settings across browser sessions"
    p.level = 1
    
    # Slide 9: Use Cases
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Use Cases"
    tf = content.text_frame
    tf.text = "• Fitness Training: Monitor exercise form and technique"
    
    p = tf.add_paragraph()
    p.text = "• Yoga Practice: Guide through pose sequences with feedback"
    p = tf.add_paragraph()
    p.text = "• Physical Therapy: Track rehabilitation exercises"
    p = tf.add_paragraph()
    p.text = "• Sports Coaching: Analyze athletic movements"
    p = tf.add_paragraph()
    p.text = "• Education: Teach proper posture and movement"
    p = tf.add_paragraph()
    p.text = "• Accessibility: Assistive technology for movement training"
    
    # Slide 10: Key Benefits
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Benefits"
    tf = content.text_frame
    tf.text = "Privacy & Security:"
    
    p = tf.add_paragraph()
    p.text = "• No data sent to external servers"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Local processing ensures privacy"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Offline functionality protects user data"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "User Experience:"
    p = tf.add_paragraph()
    p.text = "• Instant feedback with audio and visual cues"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Customizable difficulty and timing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Clear progress tracking through pose sequences"
    p.level = 1
    
    # Slide 11: Get Started Today!
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Get Started Today!"
    tf = content.text_frame
    tf.text = "Transform your training experience with Pose Recognition Web App"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nAvailable on Replit: Easy deployment and sharing"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Open Source: Customizable for your needs"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "No Installation: Run directly in your browser"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\nReady to Use:"
    p = tf.add_paragraph()
    p.text = "• Deploy instantly on Replit"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Share with your team or students"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Customize for specific use cases"
    p.level = 1
    
    # Save the presentation
    output_path = "attached_assets/Pose_Recognition_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created successfully!")
    print(f"📄 PPTX file saved: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_pose_recognition_presentation()
