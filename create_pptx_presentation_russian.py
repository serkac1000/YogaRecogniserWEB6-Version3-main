
#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_pose_recognition_presentation_russian():
    # Create presentation object
    prs = Presentation()
    
    # Set slide layout
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    bullet_slide_layout = prs.slide_layouts[1]  # Title and content
    
    # Define colors
    primary_color = RGBColor(0, 102, 204)  # Blue
    secondary_color = RGBColor(255, 153, 51)  # Orange
    
    # Slide 1: Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Веб-приложение распознавания поз"
    subtitle.text = "ИИ-система распознавания поз в реальном времени с поддержкой локальных моделей\n\nКоманда разработчиков Replit\n9 июля 2025"
    
    # Slide 2: What is pose recognition app?
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Что такое приложение распознавания поз?"
    tf = content.text_frame
    tf.text = "Распознавание поз в реальном времени с использованием веб-камеры"
    
    p = tf.add_paragraph()
    p.text = "Автономная функциональность - интернет не требуется после настройки"
    p = tf.add_paragraph()
    p.text = "Настраиваемые параметры для 1-7 поз"
    p = tf.add_paragraph()
    p.text = "Поддержка локальных моделей с интеграцией Teachable Machine"
    p = tf.add_paragraph()
    p.text = "Калибровка расстояния для оптимального позиционирования"
    p = tf.add_paragraph()
    p.text = "Звуковая обратная связь и визуальное руководство"
    
    # Slide 3: Core features
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Основные функции"
    tf = content.text_frame
    tf.text = "Система распознавания поз:"
    
    p = tf.add_paragraph()
    p.text = "• Поддержка 1-7 настраиваемых поз"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Оценка уверенности в реальном времени (порог 30-90%)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Визуальное сравнение поз с эталонными изображениями"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Отслеживание скелета с 17 ключевыми точками"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Параметры настройки:"
    p = tf.add_paragraph()
    p.text = "• Регулируемая задержка распознавания (1-10 секунд)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Пользовательские названия поз (редактируемые ярлыки)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Переключатель звукового сигнала для обратной связи об успехе"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Калибровка расстояния с руководством в реальном времени"
    p.level = 1
    
    # Slide 4: GUI settings interface
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "GUI: Интерфейс настроек"
    tf = content.text_frame
    tf.text = "Ключевые особенности:"
    
    p = tf.add_paragraph()
    p.text = "• Загрузка локальных файлов модели (model.json, metadata.json, weights.bin)"
    p = tf.add_paragraph()
    p.text = "• Флажки выбора поз для 1-7 поз"
    p = tf.add_paragraph()
    p.text = "• Загрузка эталонных изображений для каждой позы"
    p = tf.add_paragraph()
    p.text = "• Управление звуком, задержкой и порогом точности"
    p = tf.add_paragraph()
    p.text = "• Руководство по калибровке расстояния"
    
    # Try to add image if it exists
    if os.path.exists("attached_assets/GUI1_1752049448296.png"):
        slide.shapes.add_picture("attached_assets/GUI1_1752049448296.png", 
                               Inches(1), Inches(3), Inches(8), Inches(4))
    
    # Slide 5: GUI recognition interface
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "GUI: Интерфейс распознавания"
    tf = content.text_frame
    tf.text = "Функции в реальном времени:"
    tf.paragraphs[0].font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Прямая трансляция с веб-камеры с обнаружением поз"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Отображение текущей и ожидаемой позы"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Полоса уверенности с цветовым кодированием (зеленый = правильно, красный = неправильно)"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Сравнение с эталонным изображением"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Обновление и возврат к настройкам"
    p.font.size = Pt(12)
    
    # Try to add image if it exists
    if os.path.exists("attached_assets/GUI2_1752049449853.png"):
        slide.shapes.add_picture("attached_assets/GUI2_1752049449853.png", 
                               Inches(1), Inches(3), Inches(8), Inches(4))
    
    # Slide 6: Technical specifications
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Технические характеристики"
    tf = content.text_frame
    tf.text = "ИИ-фреймворк:"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• TensorFlow.js с моделями Teachable Machine"
    p.level = 1
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• PoseNet для отслеживания скелета из 17 точек"
    p.level = 1
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Локальное хранилище с IndexedDB для автономного использования"
    p.level = 1
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Совместимость:"
    p.font.size = Pt(14)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Современные браузеры с поддержкой WebRTC"
    p.level = 1
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Оптимизировано для совместимости с Windows"
    p.level = 1
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Автоматическое сжатие изображений для хранения"
    p.level = 1
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Несколько резервных разрешений камеры"
    p.level = 1
    p.font.size = Pt(12)
    
    # Slide 7: Simple setup process
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Простой процесс настройки"
    tf = content.text_frame
    tf.text = "1. Обучите модель: используйте Teachable Machine для создания модели поз"
    tf.paragraphs[0].font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "2. Загрузите файлы: импортируйте model.json, metadata.json, weights.bin"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "3. Настройте позы: выберите 1-7 поз и загрузите эталонные изображения"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "4. Настройте параметры: установите звук, задержку и предпочтения точности"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "5. Начните распознавание: запустите обнаружение поз в реальном времени"
    p.font.size = Pt(12)
    
    # Add note
    p = tf.add_paragraph()
    p.text = "\nПолностью автономно: работает полностью без интернета после первоначальной настройки"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 0, 0)
    
    # Slide 8: GUI Pose Settings Detail
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "GUI: Детальные настройки поз"
    tf = content.text_frame
    tf.text = "Настройка файлов модели:"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Загрузка model.json (файл модели)"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Загрузка metadata.json (метаданные)"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Загрузка weights.bin (веса модели)"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Выбор активных поз:"
    p.font.size = Pt(14)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Флажки для каждой позы (Class 1, 2JAIL, NOT, HLE2, LESAR, PROV, OBCHI)"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Настройка имен поз (редактируемые поля)"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Загрузка эталонных изображений для каждой позы"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Элементы управления:"
    p.font.size = Pt(14)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Звуковой сигнал (Audio Beep)"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Руководство по калибровке расстояния"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Кнопки 'Начать распознавание' и 'Начать настройку'"
    p.font.size = Pt(12)
    
    # Slide 9: Advanced features
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Расширенные функции"
    tf = content.text_frame
    tf.text = "Калибровка расстояния:"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Автоматическое определение расстояния пользователя от камеры"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Обратная связь в реальном времени для оптимального позиционирования (3-4 фута)"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Визуальные подсказки: зеленый = идеально, красный = отрегулировать расстояние"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Управление данными:"
    p.font.size = Pt(14)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Сохранение всех настроек и файлов модели локально"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Функция очистки памяти с подтверждением"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Автосохранение для выбора поз и пользовательских имен"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Постоянные настройки между сеансами браузера"
    p.font.size = Pt(12)
    
    # Slide 10: Use cases
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Случаи использования"
    tf = content.text_frame
    tf.text = "• Фитнес-тренировки: мониторинг формы и техники упражнений"
    tf.paragraphs[0].font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Практика йоги: руководство через последовательности поз с обратной связью"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Физическая терапия: отслеживание реабилитационных упражнений"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Спортивное тренерство: анализ спортивных движений"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Образование: обучение правильной осанке и движению"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Доступность: вспомогательные технологии для тренировки движений"
    p.font.size = Pt(12)
    
    # Slide 11: Key benefits
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Ключевые преимущества"
    tf = content.text_frame
    tf.text = "Конфиденциальность и безопасность:"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Никакие данные не отправляются на внешние серверы"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Локальная обработка обеспечивает конфиденциальность"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Автономная функциональность защищает пользовательские данные"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Пользовательский опыт:"
    p.font.size = Pt(14)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Мгновенная обратная связь со звуковыми и визуальными подсказками"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Настраиваемая сложность и время"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Четкое отслеживание прогресса через последовательности поз"
    p.font.size = Pt(12)
    
    # Slide 12: Get started today!
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Начните сегодня!"
    tf = content.text_frame
    tf.text = "Преобразите свой тренировочный опыт с веб-приложением распознавания поз"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nДоступно на Replit: простое развертывание и совместное использование"
    p.font.bold = True
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "Открытый исходный код: настраиваемый для ваших нужд"
    p.font.bold = True
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "Без установки: запускается прямо в вашем браузере"
    p.font.bold = True
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\nГотово к использованию:"
    p.font.size = Pt(14)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Мгновенное развертывание на Replit"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Поделитесь с вашей командой или студентами"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "• Настройте для конкретных случаев использования"
    p.font.size = Pt(12)
    
    # Save presentation
    output_path = "attached_assets/Pose_Recognition_Presentation_Russian.pptx"
    prs.save(output_path)
    print("Prezentatsiya PowerPoint uspeshno sozdana!")
    print(f"PPTX fayl sokhranyen: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_pose_recognition_presentation_russian()
